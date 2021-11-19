# -*- coding: utf-8 -*-
import sys, io, os
import pprint
import codecs
import pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.action import PP_ACTION
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob
import argparse
from PIL import Image

# japanese
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

# Index
HEIGHT = 0
WIDTH = 1
CHANNEL = 2

# Path
SRC_PPT_PATH = "src_ppt.pptx"
image_path = "./image"

# Val
ADJUST = 12700
PT = (1 / 72) * 2.54
PIXEL = 0.0264

DEFAULT_HEIGHT = 1920
DEFAULT_WIDTH = 1080

# find pptx file
def getPptName():
  name = ''
  for fl in glob.glob('./*.pptx'):
    name = os.path.split(fl)[1]
    print('Target : ' + name)

  return name

def emuToPx(emu):
  return int(emu / ADJUST * PT / PIXEL)

def getImgSize():
  img_list = glob.glob('image/*.png')
  if len(img_list) > 0:
    img = Image.open(img_list[0])
  else :
    return 1920, 1080

  return img.width, img.height

def convertMain(args):
  x, y = getImgSize()

  if args.width != 'default':
    x = int(args.width)
  if args.height != 'default':
    y = int(args.height)

  csspath = '"' + str(args.csspath) + '"'
  jqpath = '"' + str(args.jqpath) + '"'
  rwdpath = '"' + str(args.rwdpath) + '"'
  imgwidth = '"' + str(x) + '"'
  imgheight = '"' + str(y) + '"'

  # main
  print('#### Start program! ####')

  # Get ppt file name
  ppt_name = ''
  ppt_name = getPptName()
  if ppt_name == '':
    print("not found .pptx file")
    exit()

  print("image width : " + str(x) + "px")
  print("image height : " + str(y) + "px")

  # Get Ppt Infomation
  dst_ppt = Presentation(ppt_name)
  ppt_px_width = emuToPx(dst_ppt.slide_width)

  # Calc image ratio
  ratio = x / ppt_px_width

  # ppt_data = {slide_number: slide_data, ....}
  ppt_data = {}
  # ppt_name_data = {slideID: slide_number, ...}
  ppt_name_data = {}
  # ppt slide loop
  for i, sld in enumerate(dst_ppt.slides, start=1):
      # slide_data = {name: '', id: '', link_count: ''}
      ppt_slide_data = {}
      ppt_data[str(i)] = ppt_slide_data
      ppt_slide_data["name"] = 'slide' + str(i)
      ppt_slide_data["id"] = str(sld.slide_id)
      ppt_name_data[str(sld.slide_id)] = 'slide' + str(i)
      ppt_slide_data['link_count'] = 0
      # j = link count
      j = 0
      # shape loop
      for shp in sld.shapes:
        # print("shp.name : ", shp.name)
        # Not Group Shape
        if shp.shape_type != MSO_SHAPE_TYPE.GROUP:
          click_action = shp.click_action

          # check shape type
          if click_action.action == PP_ACTION.NAMED_SLIDE:
            # print("shp.name : ", shp.name)
            ppt_slide_data['link_count'] += 1
            shp_px_width = emuToPx(shp.width)
            shp_px_height = emuToPx(shp.height)
            shp_px_x1 = emuToPx(shp.left)
            shp_px_y1 = emuToPx(shp.top)
            shp_px_x2 = shp_px_x1 + shp_px_width
            shp_px_y2 = shp_px_y1 + shp_px_height
            img_px_x1 = int(shp_px_x1 * ratio)
            img_px_y1 = int(shp_px_y1 * ratio)
            img_px_x2 = int(shp_px_x2 * ratio)
            img_px_y2 = int(shp_px_y2 * ratio)
            ppt_slide_data['link' + str(j) + '_x1'] = str(img_px_x1)
            ppt_slide_data['link' + str(j) + '_y1'] = str(img_px_y1)
            ppt_slide_data['link' + str(j) + '_x2'] = str(img_px_x2)
            ppt_slide_data['link' + str(j) + '_y2'] = str(img_px_y2)
            ppt_slide_data['link' + str(j) + '_target_name'] = 'tmp'

            target = click_action.target_slide
            # DBG
            # print("target : " + str(target.slide_id))
            ppt_slide_data['link' + str(j) + '_target_id'] = target.slide_id

            j += 1
      if j == 0:
        print("Warning : Not exist link-object in slide" + str(i) + ".")

  # PPT dict loop
  print('#### Create slideX.html ####')
  for k in ppt_data:
    # html page data
    page_data = {}
    page_data['title'] = ppt_name.replace('.pptx','') + ' ' + ppt_data[k]['name']
    page_data['csspath'] = csspath
    page_data['jqpath'] = jqpath
    page_data['rwdpath'] = rwdpath
    page_data['imgpath'] = '"' + 'image/' + 'スライド' + k + '.png' + '"'
    page_data['imgwidth'] = imgwidth
    page_data['imgheight'] = imgheight

    # creat area tag
    area_all_str = ''
    if ppt_data[k]['link_count'] != 0:
      for lc in range(ppt_data[k]['link_count']):
        idx = str(lc)
        area_str = '<area  href="' + ppt_name_data[str(ppt_data[k]['link' + idx + '_target_id'])] + '.html" coords="'
        # print("LINK " + str(lc))
        area_str = area_str + str(ppt_data[k]['link' + idx + '_x1']) + ',' + str(ppt_data[k]['link' + idx + '_y1']) + ',' \
                   + str(ppt_data[k]['link' + idx + '_x2']) + ',' + str(ppt_data[k]['link' + idx + '_y2'])
        area_str = area_str + '" shape="rect">' + "\n"
        area_all_str = str(area_all_str) + area_str
    page_data['maparea'] = area_all_str

    # read temlate.html
    with open('template/template.html','r') as file:
      html = file.read()
    file.closed

    # replace {% %} to page_data
    for key, value in page_data.items():
      html = html.replace('{% ' + key + ' %}', value)

    # html output
    f = codecs.open('slide' + k + '.html', 'w', 'utf-8')
    print(html, file=f)

  print("#### Complete! ####")

def checkArguments(size):
  if size.isdigit():
    if int(size) > 0 and int(size) < 10000:
      return True
    else:
      print('Value error...')
      return False
  elif size == 'default':
    return True
  else:
    print('No other than numbers...')
    return False

def getArguments():
  parser = argparse.ArgumentParser(description='.pptxの各スライドをクリッカブルマップにコンバートします。')
  parser.add_argument('-W', '--width',  required=False, default='default' ,help='クリッカブルマップ対象の画像の幅(default:image内のpngサイズ)')
  parser.add_argument('-H', '--height',   required=False, default='default', help='クリッカブルマップ対象の画像の高さ(default:image内のpngサイズ)')
  parser.add_argument('-J', '--jqpath',  required=False, default='https://code.jquery.com/jquery-3.5.1.js', help='jQueryのパス指定(default:CDN Path)')
  parser.add_argument('-R', '--rwdpath', required=False, default='js/jquery.rwdImageMaps.js', help='jQuery RWD Image Maps のパス指定(default:js/jquery.rwdImageMaps.js)')
  parser.add_argument('-C', '--csspath', required=False, default='css/style.css', help='CSS のパス指定(default:css/style.css)')
  return parser.parse_args()

if __name__ == '__main__':
  args = getArguments()

  args_flag = False
  if checkArguments(args.width) != False and checkArguments(args.height) != False:
    args_flag = True

  if args_flag != False:
    convertMain(args)
  else:
    exit()
